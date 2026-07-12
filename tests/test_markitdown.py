from __future__ import annotations

from importlib import import_module
from importlib.metadata import version
from pathlib import Path
import tempfile
import unittest
import zipfile

from markitdown import MarkItDown
from openpyxl import Workbook
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]


class MarkItDownTest(unittest.TestCase):
    def setUp(self) -> None:
        self.converter = MarkItDown()

    def test_locked_version(self) -> None:
        self.assertEqual(version("markitdown"), "0.1.6")

    def test_representative_local_formats(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            temp = Path(temp_dir)

            docx_path = temp / "sample.docx"
            with zipfile.ZipFile(docx_path, "w") as archive:
                archive.writestr(
                    "[Content_Types].xml",
                    """<?xml version="1.0" encoding="UTF-8"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>""",
                )
                archive.writestr(
                    "_rels/.rels",
                    """<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>""",
                )
                archive.writestr(
                    "word/document.xml",
                    """<?xml version="1.0" encoding="UTF-8"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:body><w:p><w:r><w:t>DOCX conversion text</w:t></w:r></w:p></w:body>
</w:document>""",
                )

            pptx_path = temp / "sample.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            slide.shapes.title.text = "Presentation marker"
            presentation.save(pptx_path)

            xlsx_path = temp / "sample.xlsx"
            workbook = Workbook()
            worksheet = workbook.active
            worksheet["A1"] = "Spreadsheet marker"
            workbook.save(xlsx_path)

            cases = (
                (docx_path, "DOCX conversion text"),
                (pptx_path, "Presentation marker"),
                (xlsx_path, "Spreadsheet marker"),
            )
            for path, marker in cases:
                with self.subTest(path=path.name):
                    result = self.converter.convert(path)
                    self.assertIn(marker, result.text_content)

    def test_pdf_fixture(self) -> None:
        fixture = (
            ROOT
            / ".agents/skills/source-extraction/references/providers/mineru/tests/fixtures"
            / "1512.03385-resnet.pdf"
        )
        result = self.converter.convert(fixture)
        self.assertIn("Residual", result.text_content)
        self.assertIn("Learning", result.text_content)

    def test_retained_extras_import(self) -> None:
        modules = (
            "azure.ai.documentintelligence",
            "azure.identity",
            "olefile",
            "pydub",
            "speech_recognition",
            "youtube_transcript_api",
        )
        for module_name in modules:
            with self.subTest(module=module_name):
                self.assertIsNotNone(import_module(module_name))


if __name__ == "__main__":
    unittest.main()
